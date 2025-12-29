from pptx import Presentation
from pptx.util import Inches
import os
import sys

from pptx.util import Pt
from pptx.dml.color import RGBColor

# Adjust path to include src if running from root
sys.path.append(os.path.join(os.path.dirname(__file__), '../../'))

from src.schema.slide_schema import PresentationDeck, SlideContent, SlideType
from src.schema.slide_schema import PresentationDeck, SlideContent, SlideType
from src.config.style_config import StyleConfig, TextStyle, JMDCFont, JMDCColor
from src.utils.chart_builder import ChartBuilder

class SlideBuilder:
    def __init__(self, template_path: str):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        self.prs = Presentation(template_path)
        
        # Hardcoded layout map based on analysis (could be dynamic later)
        # Template structure:
        # Master 0 (Index 0): Title, TOC, Section, Content, Back
        self.layout_map = {
            SlideType.COVER: 0,
            SlideType.TOC: 1,
            SlideType.SECTION: 2,
            SlideType.CONTENT: 3,
            SlideType.BACK_COVER: 4
        }

    def _apply_text_style(self, shape, style_type: TextStyle, text: str):
        """
        Applies strict styling to a shape's text frame.
        """
        if not shape.has_text_frame:
            return
            
        text_frame = shape.text_frame
        text_frame.clear() # Clear existing dummy text/formatting
        
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        
        style = StyleConfig.get_font_style(style_type)
        
        font = run.font
        font.name = style["name"]
        font.size = Pt(style["size"]) if style["size"] else None
        font.bold = style["bold"]
        font.italic = style["italic"]
        
        if style["color"]:
            font.color.rgb = style["color"]
            
        # Support CJK fonts explicitly if needed (some libraries need extra handling for Asian fonts)
        # python-pptx handles `font.name` reasonably well for installed fonts.


    def build(self, deck: PresentationDeck, output_path: str):
        """
        Generates the presentation and saves it to output_path.
        """
        # Clear existing slides (optional, but usually we want a fresh start from template masters)
        # Note: python-pptx doesn't easily allow deleting all slides while keeping masters unless we start with a clean template.
        # Assuming the input template is a clean "potx" style or we just append.
        # For this logic, we append. If the template has dummy slides, they will remain at the beginning.
        # TODO: Implement slide removal if needed.

        for slide_content in deck.slides:
            self._create_slide(slide_content)
        
        self.prs.save(output_path)
        print(f"Presentation saved to {output_path}")

    def _create_slide(self, content: SlideContent):
        layout_index = self.layout_map.get(content.type, 3) # Default to content
        layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(layout)
        
        # Map content to placeholders
        # This mapping is specific to the JMDC template structure we analyzed
        
        # Data Mapping
        
        # Title mapping
        if content.title:
            if slide.shapes.title:
                # Determine title style based on slide type
                t_style = TextStyle.TITLE_MAIN if content.type == SlideType.COVER else TextStyle.TITLE_SLIDE
                self._apply_text_style(slide.shapes.title, t_style, content.title)
        
        # Subtitle / Body mapping
        
        if content.type == SlideType.COVER:
            if content.subtitle:
                try:
                    self._apply_text_style(slide.placeholders[1], TextStyle.SUBTITLE, content.subtitle)
                except:
                    pass
        
        elif content.type == SlideType.SECTION:
            if content.subtitle:
                try:
                    self._apply_text_style(slide.placeholders[1], TextStyle.SUBTITLE, content.subtitle)
                except:
                    pass

        elif content.type == SlideType.CONTENT:
            if content.body:
                try:
                    # Prefer the object placeholder for general content
                    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
                        self._apply_text_style(slide.placeholders[1], TextStyle.BODY, content.body)
                except:
                    pass
            
            if content.subtitle:
                try:
                    # Use the top text placeholder for subtitle
                    if 13 in [p.placeholder_format.idx for p in slide.placeholders]:
                        self._apply_text_style(slide.placeholders[13], TextStyle.SUBTITLE, content.subtitle)
                except:
                    pass



        # Chart Handling
        if content.chart:
            try:
                # Find a suitable location for the chart.
                # Default to the center "Object" placeholder (idx 1 usually) if available,
                # otherwise use a calculated center position.
                
                # Check for Content Placeholder (idx 1)
                target_ph = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        target_ph = ph
                        break
                
                if target_ph:
                    x, y = target_ph.left, target_ph.top
                    cx, cy = target_ph.width, target_ph.height
                else:
                    # Fallback to roughly center
                    x, y = Inches(1), Inches(2)
                    cx, cy = Inches(8), Inches(4.5)

                ChartBuilder.create_chart(slide, content.chart, x, y, cx, cy)
                
            except Exception as e:
                print(f"Warning: Failed to create chart: {e}")

        # Advanced Elements Handling (PDF Reconstruction)
        if content.elements:
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
            
            for elem in content.elements:
                if not elem.rect:
                    continue
                    
                # Convert normalized rect to EMU/Inches
                x = int(elem.rect[0] * slide_width)
                y = int(elem.rect[1] * slide_height)
                w = int(elem.rect[2] * slide_width)
                h = int(elem.rect[3] * slide_height)
                
                if elem.type == "text":
                    # Create Text Box
                    textbox = slide.shapes.add_textbox(x, y, w, h)
                    tf = textbox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = elem.content
                    
                    # Apply Styles
                    run.font.name = JMDCFont.REGULAR
                    
                    # Heuristic font size mapping or direct use
                    if elem.font_size:
                        # Map PDF point size roughly to PPTX point size
                        # Often PDF points are similar.
                        # Enforce minimum size for readability?
                        run.font.size = Pt(elem.font_size)
                    else:
                        run.font.size = Pt(12) 
                        
                    # Color? Default to Black or Dark Grey
                    run.font.color.rgb = JMDCColor.TEXT_MAIN

                elif elem.type == "image":
                    # Add Picture
                    if os.path.exists(elem.content):
                        try:
                            slide.shapes.add_picture(elem.content, x, y, w, h)
                        except Exception as e:
                            print(f"Warning: Failed to add image {elem.content}: {e}")
