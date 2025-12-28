from pptx.chart.data import CategoryChartData, ChartData as PptxChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from src.schema.slide_schema import ChartData, ChartType

class ChartBuilder:
    @staticmethod
    def create_chart(slide, chart_def: ChartData, x, y, cx, cy):
        """
        Creates a chart on the given slide at specified coordinates.
        """
        chart_data = CategoryChartData()
        chart_data.categories = chart_def.categories

        # Add series
        for name, values in chart_def.series.items():
            chart_data.add_series(name, values)

        # Determine Chart Type
        xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED # Default
        if chart_def.type == ChartType.BAR_CLUSTERED:
            xl_type = XL_CHART_TYPE.BAR_CLUSTERED
        elif chart_def.type == ChartType.LINE:
            xl_type = XL_CHART_TYPE.LINE
        elif chart_def.type == ChartType.PIE:
            xl_type = XL_CHART_TYPE.PIE
            
        # Add Chart
        graphic_frame = slide.shapes.add_chart(
            xl_type, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        
        # Basic Styling
        if chart_def.title:
            chart.chart_title.text_frame.text = chart_def.title
        
        # Tweak fonts for legibility (using rough defaults for now)
        # In a real scenario, we'd use StyleConfig here too
        if chart.has_legend:
            chart.legend.font.size = Pt(10)
            chart.legend.include_in_layout = False

        return chart
