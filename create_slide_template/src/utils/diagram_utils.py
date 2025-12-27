from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_diagram(slide, diagram_info, colors=None):
    """
    スライドに図解（シェイプ群）を追加する
    diagram_info: {
        "type": "PROCESS" | "CYCLE" | "LIST",
        "items": ["Item 1", "Item 2", ...],
        "position": {...}
    }
    """
    dtype = diagram_info.get("type", "PROCESS").upper()
    items = diagram_info.get("items", [])
    if not items:
        return

    pos = diagram_info.get("position", {})
    base_left = Inches(pos.get("left", 1.0))
    base_top = Inches(pos.get("top", 2.0))
    base_width = Inches(pos.get("width", 8.0))
    base_height = Inches(pos.get("height", 3.0))

    if dtype == "PROCESS":
        _create_process_flow(slide, items, base_left, base_top, base_width, base_height, colors)
    elif dtype == "CYCLE":
        _create_cycle(slide, items, base_left, base_top, base_width, base_height, colors)
    else:
        # Default to simple list boxes
        _create_process_flow(slide, items, base_left, base_top, base_width, base_height, colors)

def _create_process_flow(slide, items, left, top, width, height, colors):
    """左から右への矢印付きプロセスフロー"""
    count = len(items)
    # マージンとシェイプの計算
    margin = Inches(0.2)
    shape_width = (width - (margin * (count - 1))) / count
    shape_height = Inches(1.0)
    
    current_left = left
    
    for i, text in enumerate(items):
        # 図形の追加 (角丸四角形)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            current_left, top, shape_width, shape_height
        )
        shape.text = text
        
        # 色設定
        if colors:
            color_hex = colors[i % len(colors)].lstrip('#')
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            
        # 矢印の追加 (最後以外)
        if i < count - 1:
            arrow_left = current_left + shape_width
            arrow_top = top + (shape_height / 2) - Inches(0.1) # 中心
            arrow_w = margin
            # 矢印シェイプは面倒なので、簡単な「右矢印」図形を挟むか、コネクタを使う
            # ここではコネクタ（線）ではなく、右矢印図形を小さく配置する
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left, top + (shape_height/2) - Inches(0.15), 
                margin, Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor.from_string("888888") # Grey arrow

        current_left += shape_width + margin

def _create_cycle(slide, items, left, top, width, height, colors):
    """簡易的なサイクル図（円形配置）"""
    import math
    count = len(items)
    center_x = left + (width / 2)
    center_y = top + (height / 2)
    radius = min(width, height) / 3
    
    shape_w = Inches(1.5)
    shape_h = Inches(0.8)
    
    for i, text in enumerate(items):
        angle = (2 * math.pi * i) / count - (math.pi / 2) # 上から開始
        
        # 中心からのオフセット
        x = center_x + (radius * math.cos(angle)) - (shape_w / 2)
        y = center_y + (radius * math.sin(angle)) - (shape_h / 2)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y, shape_w, shape_h
        )
        shape.text = text
        
        if colors:
            color_hex = colors[i % len(colors)].lstrip('#')
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            
        # 矢印（次の要素へ）- 複雑になるので省略、または中心から外への線など
        # ここではシンプルに円配置のみとする
