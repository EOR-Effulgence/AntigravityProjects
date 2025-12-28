import json
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

def get_solid_color(params):
    """
    Helper to extract color from a fill/font object.
    Returns hex string or theme color name.
    """
    if not hasattr(params, 'color'):
        return None
    try:
        color = params.color
        if color.type == MSO_THEME_COLOR.NOT_THEME_COLOR:
            if color.rgb:
                return str(color.rgb)
        else:
            return f"THEME_{color.theme_color}"
    except:
        return None
    return None

def analyze_shape_text_style(shape):
    """
    Extracts text styles (font, size, color, etc.) from a shape's text frame.
    Checks the first paragraph/run as a representative sample.
    """
    style = {}
    if not shape.has_text_frame:
        return None
        
    tf = shape.text_frame
    
    # Text Frame level (margins, alignment could be here)
    # style["margin_left"] = tf.margin_left
    
    # Paragraph level (check first paragraph)
    if tf.paragraphs:
        p = tf.paragraphs[0]
        style["alignment"] = str(p.alignment)
        # style["line_spacing"] = p.line_spacing
        # style["space_before"] = p.space_before
        # style["space_after"] = p.space_after
        
        # Run level (font properties)
        if p.runs:
            r = p.runs[0]
            font = r.font
            style["font_name"] = font.name
            style["font_size"] = font.size.pt if font.size else None
            style["font_bold"] = font.bold
            style["font_italic"] = font.italic
            style["font_color"] = get_solid_color(font)
            
    return style

def analyze_presentation(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return

    prs = Presentation(pptx_path)
    
    info = {
        "slide_width": prs.slide_width.inches,
        "slide_height": prs.slide_height.inches,
        "masters": []
    }

    print(f"Analyzing {os.path.basename(pptx_path)}...")

    for i, master in enumerate(prs.slide_masters):
        master_data = {
            "id": f"master_{i}",
            "name": master.name,
            "layouts": []
        }
        
        print(f"\nMaster {i+1}: {master.name}")
        
        for layout in master.slide_layouts:
            try:
                layout_id = layout.slide_id
            except AttributeError:
                layout_id = "unknown"

            layout_data = {
                "index": prs.slide_layouts.index(layout) if layout in prs.slide_layouts else -1,
                "id": layout_id,
                "name": layout.name,
                "placeholders": []
            }
            
            # Identify key visual elements
            for shape in layout.placeholders:
                ph_type = shape.placeholder_format.type
                
                ph_data = {
                    "idx": shape.placeholder_format.idx,
                    "name": shape.name,
                    "type": str(ph_type),
                    "pos": {
                        "left": shape.left.inches,
                        "top": shape.top.inches,
                        "width": shape.width.inches,
                        "height": shape.height.inches
                    },
                    "style": analyze_shape_text_style(shape)
                }
                layout_data["placeholders"].append(ph_data)
            
            master_data["layouts"].append(layout_data)
        
        info["masters"].append(master_data)

    # Analyze actual slides to find concrete style examples
    info["slides"] = []
    print(f"\nAnalyzing {len(prs.slides)} slides...")
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "index": i,
            "shapes": []
        }
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                try:
                    style = analyze_shape_text_style(shape)
                    if style and style.get("font_name"): # Only interesting if it has font info
                        shape_data = {
                            "name": shape.name,
                            "text_preview": shape.text[:50], # First 50 chars
                            "style": style
                        }
                        slide_data["shapes"].append(shape_data)
                except Exception as e:
                    pass
        
        if slide_data["shapes"]:
            info["slides"].append(slide_data)

    output_path = "template_style_analysis.json"
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(info, f, indent=2, ensure_ascii=False)
    
    print(f"\nAnalysis complete. Structure saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python analyze_pptx.py <path_to_pptx>")
    else:
        analyze_presentation(sys.argv[1])
