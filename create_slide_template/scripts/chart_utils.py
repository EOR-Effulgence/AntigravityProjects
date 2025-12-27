import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# チャートタイプのマッピング
CHART_TYPES = {
    "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
    "LINE": XL_CHART_TYPE.LINE,
    "LINE_MARKERS": XL_CHART_TYPE.LINE_MARKERS,
    "PIE": XL_CHART_TYPE.PIE,
    "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
    "AREA": XL_CHART_TYPE.AREA,
    "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "BAR": XL_CHART_TYPE.BAR_CLUSTERED,
    "STACKED_COLUMN": XL_CHART_TYPE.COLUMN_STACKED,
}

def create_chart_data(data_config):
    """
    JSONデータからChartDataオブジェクトを作成する
    data_config format:
    {
        "categories": ["cat1", "cat2", ...],
        "series": [
            {"name": "series1", "values": [val1, val2, ...]},
            ...
        ]
    }
    """
    chart_data = CategoryChartData()
    
    # カテゴリ設定
    categories = data_config.get("categories", [])
    chart_data.categories = categories
    
    # シリーズ設定
    series_list = data_config.get("series", [])
    for series in series_list:
        name = series.get("name", "")
        values = series.get("values", [])
        
        # 値の検証と変換
        float_values = []
        for v in values:
            try:
                float_values.append(float(v))
            except (ValueError, TypeError):
                float_values.append(0.0)
                
        chart_data.add_series(name, float_values)
        
    return chart_data

def parse_chart_data_from_csv(csv_text: str) -> dict:
    """
    Markdownのチャートブロック用CSVテキストを解析して辞書形式に変換する
    
    Expected format:
    Title (Optional, handled outside or first line)
    Header1, Header2, Header3
    Label1, Val1, Val2
    Label2, Val3, Val4
    
    Returns:
    {
        "categories": ["Label1", "Label2"],
        "series": [
            {"name": "Header2", "values": [Val1, Val3]},
            {"name": "Header3", "values": [Val2, Val4]}
        ]
    }
    """
    import csv
    import io
    
    lines = csv_text.strip().split('\n')
    if not lines:
        return {}
        
    # タイトル行の扱い（オプション）
    # ここでは純粋なCSVパートのみを受け取る前提とするが、
    # もし1行だけ列数が違うなどが検知できればタイトルとして扱うロジックも追加可能
    
    reader = csv.reader(lines)
    header = next(reader, None)
    
    if not header:
        return {}
        
    # Header: [CategoryKey, SeriesName1, SeriesName2, ...]
    if len(header) < 2:
        return {}
        
    series_names = header[1:]
    categories = []
    series_values = [[] for _ in series_names]
    
    for row in reader:
        if not row: continue
        # 行の長さが足りない場合は空文字で埋める
        if len(row) < len(header):
            row += [''] * (len(header) - len(row))
            
        categories.append(row[0])
        
        for i, val_str in enumerate(row[1:]):
            if i < len(series_values):
                try:
                    # 数値変換（カンマ除去など）
                    clean_val = val_str.replace(',', '').strip()
                    val = float(clean_val)
                except ValueError:
                    val = 0.0
                series_values[i].append(val)
                
    # 構築
    series_list = []
    for i, name in enumerate(series_names):
        series_list.append({
            "name": name.strip(),
            "values": series_values[i]
        })
        
    return {
        "categories": categories,
        "series": series_list
    }


def add_chart_to_slide(slide, chart_info):
    """
    スライドにグラフを追加する
    chart_info format:
    {
        "type": "COLUMN_CLUSTERED",
        "title": "Chart Title",
        "data": { ... },
        "position": {"left": 1.0, "top": 2.0, "width": 6.0, "height": 4.0}
    }
    """
    chart_type_str = chart_info.get("type", "COLUMN_CLUSTERED")
    chart_type = CHART_TYPES.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    data_config = chart_info.get("data", {})
    chart_data = create_chart_data(data_config)
    
    # 位置情報 (インチ単位)
    pos = chart_info.get("position", {})
    left = Inches(pos.get("left", 1.0))
    top = Inches(pos.get("top", 2.0))
    width = Inches(pos.get("width", 6.0))
    height = Inches(pos.get("height", 4.0))
    
    # グラフを追加
    graphic_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    
    # タイトル設定
    title = chart_info.get("title")
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        
    # 凡例設定 (デフォルトで表示)
    chart.has_legend = True
    chart.legend.include_in_layout = False

def add_table_to_slide(slide, table_info):
    """
    スライドに表を追加する
    table_info format:
    {
        "data": [
            ["Header1", "Header2"],
            ["Row1Col1", "Row1Col2"]
        ],
        "position": {"left": 1.0, "top": 2.0, "width": 6.0, "height": 2.0}
    }
    """
    data = table_info.get("data", [])
    if not data:
        return
        
    rows = len(data)
    cols = len(data[0])
    
    # 位置情報
    pos = table_info.get("position", {})
    left = Inches(pos.get("left", 1.0))
    top = Inches(pos.get("top", 2.0))
    width = Inches(pos.get("width", 6.0))
    height = Inches(pos.get("height", 2.0))
    
    # 表を追加
    graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = graphic_frame.table
    
    # データを設定
    for i, row_data in enumerate(data):
        for j, cell_value in enumerate(row_data):
            if j < cols:
                cell = table.cell(i, j)
                if cell_value is not None:
                    cell.text = str(cell_value)
                    
                # ヘッダー行のスタイル (簡易設定)
                if i == 0:
                    pass # ここでスタイル設定可能
