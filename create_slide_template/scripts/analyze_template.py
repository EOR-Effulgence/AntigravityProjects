#!/usr/bin/env python3
"""
JMDCテンプレートを詳細に解析してレイアウト・プレースホルダー情報を抽出するスクリプト
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import PP_PLACEHOLDER


def analyze_template(template_path: str, output_path: str = None):
    """テンプレートを解析してレイアウト情報を抽出"""
    prs = Presentation(template_path)
    
    result = {
        "template_file": str(template_path),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_width_inches": prs.slide_width / 914400,
        "slide_height_inches": prs.slide_height / 914400,
        "slide_masters": [],
        "existing_slides": []
    }
    
    # スライドマスターとレイアウトを解析
    for master_idx, master in enumerate(prs.slide_masters):
        master_info = {
            "index": master_idx,
            "name": master.name if hasattr(master, 'name') else f"Master {master_idx}",
            "layouts": []
        }
        
        # レイアウトを解析
        for layout_idx, layout in enumerate(master.slide_layouts):
            layout_info = {
                "index": layout_idx,
                "name": layout.name,
                "placeholders": []
            }
            
            # プレースホルダーを解析
            # プレースホルダーのインデックス順ではなく、視覚的な順序（上から下、左から右）でソートしたいが
            # ここではまず生の情報を全て取得する
            for ph in layout.placeholders:
                try:
                    name = ph.name
                except AttributeError:
                    name = "unknown"
                
                ph_info = {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": name,
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                    "left_inch": ph.left / 914400 if ph.left else 0.0,
                    "top_inch": ph.top / 914400 if ph.top else 0.0,
                    "width_inch": ph.width / 914400 if ph.width else 0.0,
                    "height_inch": ph.height / 914400 if ph.height else 0.0
                }
                layout_info["placeholders"].append(ph_info)
            
            # プレースホルダーをTop順、Left順にソートして出力（視覚的理解のため）
            layout_info["placeholders"].sort(key=lambda x: (x['top'], x['left']))
            
            master_info["layouts"].append(layout_info)
        
        result["slide_masters"].append(master_info)
    
    # 結果を出力
    if output_path:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2, default=str)
        print(f"分析結果を保存しました: {output_path}")
    
    return result


def print_summary(result: dict):
    """分析結果のサマリーを表示"""
    print("\n" + "="*80)
    print("テンプレート詳細分析サマリー")
    print("="*80)
    print(f"\nファイル: {result['template_file']}")
    print(f"スライドサイズ: {result['slide_width_inches']:.2f}\" x {result['slide_height_inches']:.2f}\"")
    
    for master in result["slide_masters"]:
        # 必要なレイアウト（表紙、目次、中見出し、コンテンツ）のみ詳細表示
        target_layouts = ["表紙", "目次", "中見出し", "コンテンツ", "裏表紙"]
        
        for layout in master["layouts"]:
            if any(t in layout['name'] for t in target_layouts):
                print(f"\n  [{layout['index']}] {layout['name']}")
                print("-" * 60)
                print(f"  {'idx':<4} | {'type':<20} | {'name':<20} | {'pos (inch)':<15} | {'size (inch)':<15}")
                print("-" * 60)
                
                for ph in layout["placeholders"]:
                    pos = f"({ph['left_inch']:.1f}, {ph['top_inch']:.1f})"
                    size = f"{ph['width_inch']:.1f} x {ph['height_inch']:.1f}"
                    ph_type = ph['type'].replace('PP_PLACEHOLDER.', '') if 'PP_PLACEHOLDER.' in ph['type'] else ph['type']
                    
                    print(f"  {ph['idx']:<4} | {ph_type[:20]:<20} | {ph['name'][:20]:<20} | {pos:<15} | {size:<15}")


if __name__ == "__main__":
    import sys
    
    # デフォルトのテンプレートパス
    script_dir = Path(__file__).parent.parent
    default_template = script_dir / "JMDC2022_16対9(標準)_v1.1.pptx"
    output_json = script_dir / "template_analysis_detail.json"
    
    template_path = sys.argv[1] if len(sys.argv) > 1 else str(default_template)
    
    print(f"詳細解析中: {template_path}")
    result = analyze_template(template_path, str(output_json))
    print_summary(result)
