#!/usr/bin/env python3
"""
JMDCテンプレートを使用してPowerPointスライドを生成するスクリプト
"""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def load_slide_data(json_path: str) -> dict:
    """JSONファイルからスライドデータを読み込む"""
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


# グラフ・表・画像ユーティリティのインポート
scripts_dir = Path(__file__).parent
sys.path.append(str(scripts_dir))

chart_utils = None
image_utils = None

try:
    import chart_utils
except ImportError:
    print("警告: chart_utils.py が見つかりません。グラフ・表の生成はスキップされます。")

try:
    import image_utils
except ImportError:
    print("警告: image_utils.py が見つかりません。画像生成はスキップされます。")


def load_mapping(config_path: str) -> dict:
    """レイアウトマッピング設定を読み込む"""
    with open(config_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def create_presentation(template_path: str, slide_data: dict, output_path: str, mapping_path: str):
    """
    テンプレートを使用してプレゼンテーションを作成
    """
    prs = Presentation(template_path)
    mapping_config = load_mapping(mapping_path)
    layout_mapping = mapping_config["layouts"]
    
    # テンプレートのレイアウト名とインデックスのマッピング
    ppt_layout_indices = {}
    for idx, layout in enumerate(prs.slide_masters[0].slide_layouts):
        ppt_layout_indices[layout.name] = idx
    
    print(f"テンプレート内のレイアウト: {list(ppt_layout_indices.keys())}")
    
    # 既存のサンプルスライドを削除
    if slide_data.get("clear_existing", False):
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    
    # 新しいスライドを追加
    for i, slide_info in enumerate(slide_data.get("slides", [])):
        layout_name = slide_info.get("layout", "コンテンツ")
        
        # マッピング定義にあるレイアウトか確認
        if layout_name not in layout_mapping:
            print(f"警告: マッピング定義にないレイアウト '{layout_name}' です。スキップします。")
            continue
            
        # テンプレートに存在するレイアウトか確認
        if layout_name not in ppt_layout_indices:
            print(f"警告: テンプレートにレイアウト '{layout_name}' が見つかりません。")
            continue
        
        # スライド追加
        layout_idx = ppt_layout_indices[layout_name]
        layout = prs.slide_masters[0].slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        # マッピング定義に基づいてコンテンツを設定
        ph_map = layout_mapping[layout_name]
        
        for field, config in ph_map.items():
            if field == "description": continue
            
            target_idx = config["idx"]
            
            # 入力データにそのフィールドの値があるか確認
            content_value = slide_info.get(field)
            if content_value is None:
                if field == "content" and "body" in slide_info:
                    content_value = slide_info["body"]
                elif field == "body" and "content" in slide_info:
                    content_value = slide_info["content"]
            
            if content_value:
                # 該当するインデックスのプレースホルダーを探す
                try:
                    shape = slide.placeholders[target_idx]
                    
                    if isinstance(content_value, list):
                        # リスト形式のコンテンツ
                        tf = shape.text_frame
                        tf.clear()
                        for i, item in enumerate(content_value):
                            if i == 0:
                                p = tf.paragraphs[0]
                                p.text = item
                            else:
                                p = tf.add_paragraph()
                                p.text = item
                                p.level = 0
                            
                            # フォント設定
                            if p.runs:
                                for run in p.runs:
                                    run.font.name = "Noto Sans CJK JP"
                                    # run.font.language_id = MSO_LANGUAGE_ID.JAPANESE # 必要なら
                            
                            # 段落全体のフォント設定（念のため）
                            p.font.name = "Noto Sans CJK JP"

                    else:
                        # テキスト形式
                        shape.text = str(content_value)
                        # テキストフレーム全体のフォントを設定
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.name = "Noto Sans CJK JP"
                            for run in paragraph.runs:
                                run.font.name = "Noto Sans CJK JP"
                        
                except KeyError:
                    print(f"警告: スライド[{i}] '{layout_name}' のプレースホルダー idx={target_idx} ({field}) が見つかりません。")
        
        # グラフの追加
        if "charts" in slide_info and chart_utils:
            for chart_info in slide_info["charts"]:
                try:
                    chart_utils.add_chart_to_slide(slide, chart_info)
                    print(f"  - グラフ追加: {chart_info.get('title', 'No Title')}")
                except Exception as e:
                    print(f"  - グラフ追加エラー: {e}")

        # 表の追加
        if "tables" in slide_info and chart_utils:
            for table_info in slide_info["tables"]:
                try:
                    chart_utils.add_table_to_slide(slide, table_info)
                    print(f"  - 表追加")
                except Exception as e:
                    print(f"  - 表追加エラー: {e}")

        # 画像の追加 (NanoBanaNana連携)
        if "images" in slide_info and image_utils:
            for image_info in slide_info["images"]:
                try:
                    image_utils.add_image_to_slide(slide, image_info)
                    print(f"  - 画像追加: {image_info.get('prompt', 'No Prompt')[:20]}...")
                except Exception as e:
                    print(f"  - 画像追加エラー: {e}")

        print(f"スライド追加: {layout_name} - {slide_info.get('title', 'No Title')}")
    
    # 保存
    prs.save(output_path)
    print(f"\nプレゼンテーションを保存しました: {output_path}")


def main():
    if len(sys.argv) < 3:
        print("使用方法: python generate_slide.py <input.json> <output.pptx>")
        sys.exit(1)
    
    input_json = sys.argv[1]
    output_pptx = sys.argv[2]

    # 出力パスの自動調整（日付フォルダ分け）
    output_path_obj = Path(output_pptx)
    if len(output_path_obj.parts) == 1: # ファイル名のみの場合
        import datetime
        today_str = datetime.datetime.now().strftime('%Y%m%d')
        output_dir = Path("output") / today_str
        output_dir.mkdir(parents=True, exist_ok=True)
        output_pptx = str(output_dir / output_pptx)
        print(f"出力先フォルダを自動設定しました: {output_pptx}")
    else:
        # ディレクトリ指定がある場合はその親ディレクトリを作成しておく
        output_path_obj.parent.mkdir(parents=True, exist_ok=True)
    
    # パス設定
    script_dir = Path(__file__).parent.parent
    template_path = script_dir / "JMDC2022_16対9(標準)_v1.1.pptx"
    mapping_path = script_dir / "config" / "layout_mapping.json"
    
    slide_data = load_slide_data(input_json)
    create_presentation(str(template_path), slide_data, output_pptx, str(mapping_path))


if __name__ == "__main__":
    main()
